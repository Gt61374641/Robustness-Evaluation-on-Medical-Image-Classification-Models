"""Generate a minimal 7-slide group-meeting deck (English).

Usage:  python scripts/make_slides.py
Output: meeting_slides.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "figures" / "combined"

NAVY = RGBColor(0x1F, 0x2D, 0x3D)
GRAY = RGBColor(0x44, 0x4A, 0x52)
ACCENT = RGBColor(0x2C, 0x6E, 0x9E)
LIGHT = RGBColor(0x8A, 0x90, 0x99)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    return prs.slides.add_slide(BLANK)


def textbox(s, left, top, width, height):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def set_run(p, text, size, color, bold=False, italic=False):
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.color.rgb = color; f.bold = bold; f.italic = italic
    f.name = "Calibri"
    return r


def title(s, text, sub=None):
    tf = textbox(s, 0.7, 0.5, 12, 1.1)
    p = tf.paragraphs[0]; set_run(p, text, 30, NAVY, bold=True)
    # thin accent rule (no shadow -> avoids ghosting)
    line = s.shapes.add_shape(1, Inches(0.75), Inches(1.45), Inches(2.2), Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT; line.line.fill.background()
    line.shadow.inherit = False
    if sub:
        tfs = textbox(s, 0.75, 1.55, 12, 0.6)
        set_run(tfs.paragraphs[0], sub, 15, LIGHT, italic=True)


def bullets(s, items, left=0.85, top=1.9, width=11.7, size=18, gap=10):
    tf = textbox(s, left, top, width, 5)
    for i, (txt, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.level = lvl
        prefix = "•  " if lvl == 0 else "–  "
        set_run(p, prefix, size, ACCENT if lvl == 0 else LIGHT, bold=(lvl == 0))
        set_run(p, txt, size if lvl == 0 else size - 2, GRAY if lvl == 0 else LIGHT)


def image(s, path, top, height=None, width=None):
    kw = {}
    if width: kw["width"] = Inches(width)
    if height: kw["height"] = Inches(height)
    pic = s.shapes.add_picture(str(path), Inches(0), Inches(top), **kw)
    pic.left = int((SW - pic.width) / 2)  # center horizontally
    return pic


def footer(s, n):
    tf = textbox(s, 11.6, 7.0, 1.6, 0.4)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    set_run(p, f"{n} / 7", 11, LIGHT)


# ---- Slide 1: title ----
s = slide()
tf = textbox(s, 1.0, 2.5, 11.3, 2.0)
set_run(tf.paragraphs[0], "Robustness Evaluation on Medical Image Classification Models", 36, NAVY, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(14)
set_run(p, "Model complexity & adversarial training across three medical datasets", 20, ACCENT)
tf2 = textbox(s, 1.0, 5.4, 11.3, 1.0)
set_run(tf2.paragraphs[0], "Two-week progress report   ·   23 June 2026", 16, LIGHT)
bar = s.shapes.add_shape(1, Inches(1.05), Inches(4.55), Inches(3.0), Pt(4))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
bar.shadow.inherit = False

# ---- Slide 2: background & objectives ----
s = slide(); title(s, "Background & Objectives")
bullets(s, [
    ("Deep models are highly vulnerable to imperceptible adversarial perturbations — a real safety risk for medical diagnosis.", 0),
    ("We study how model complexity affects adversarial robustness, following Rodriguez et al. (2022, BMC).", 0),
    ("H1 — Under standard training, how does model complexity affect robustness?", 0),
    ("H2 — Does adversarial training (AT) change the complexity–robustness relationship?", 0),
])
footer(s, 2)

# ---- Slide 3: setup ----
s = slide(); title(s, "Experimental Setup")
bullets(s, [
    ("Models — ResNet-18 / 34 / 50 / 101 / 152 (complexity ladder, ImageNet-pretrained).", 0),
    ("Datasets — Chest X-ray (primary), Malaria (patient-level split), OCT2017 (4-class).", 0),
    ("Attacks — FGSM + PGD ε-sweep; strong evaluation: PGD-50 + 5 random restarts.", 0),
    ("Defense — PGD adversarial training (AT) on representative models.", 0),
    ("Clean baselines — Chest 0.83–0.85 · Malaria 0.97 · OCT 0.99.", 0),
])
footer(s, 3)

# ---- Slide 4: H1 ----
s = slide(); title(s, "H1 — Complexity–robustness is non-monotonic",
                   "Top: robustness vs attack budget (ε)   ·   Bottom: robustness vs model complexity (U-shape)")
image(s, FIG / "H1_dual_view.png", top=2.0, width=10.0)
footer(s, 4)

# ---- Slide 5: H2 ----
s = slide(); title(s, "H2 — Adversarial training restores robustness", "Larger / mid-capacity models benefit more")
image(s, FIG / "H2_at_across_datasets.png", top=2.1, width=11.6)
bullets(s, [
    ("Standard models ≈ 0 at 8/255 under strong PGD; AT lifts R50/R152 to 0.5–0.9 (Malaria 0.90 @8/255).", 0),
    ("Smallest model (R18) is hardest to converge under AT — consistent with “complexity helps AT”.", 0),
], top=5.8, size=14, gap=5)
footer(s, 5)

# ---- Slide 6: methodology ----
s = slide(); title(s, "Methodology Improvements")
bullets(s, [
    ("Reliable PGD evaluation (PGD-50 + restarts) with an automatic large-ε class-collapse diagnostic.", 0),
    ("Harmonized AT loop — same loss / scheduler / AMP / checkpoint-selection as standard training.", 0),
    ("AT stabilization — ε and learning-rate warmup to prevent early-training collapse.", 0),
    ("Fair, leak-free protocol — fixed test subset, stratified splits, patient-level Malaria split, resume-integrity guard.", 0),
])
footer(s, 6)

# ---- Slide 7: status & next ----
s = slide(); title(s, "Status & Next Steps")
bullets(s, [
    ("Done", 0),
    ("Standard training + robustness sweep on all three datasets.", 1),
    ("Adversarial training on representative models; cross-dataset figures.", 1),
    ("In progress / next", 0),
    ("Finish AT convergence for R18 / OCT (warmup re-runs).", 1),
    ("Interpretability — Grad-CAM & decision-boundary visualizations.", 1),
    ("Write Results & Discussion.", 1),
])
footer(s, 7)

out = ROOT / "meeting_slides.pptx"
prs.save(str(out))
print(f"wrote {out}")
